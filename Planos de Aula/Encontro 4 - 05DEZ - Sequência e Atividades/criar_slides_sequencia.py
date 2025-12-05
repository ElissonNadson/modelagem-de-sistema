# -*- coding: utf-8 -*-
"""
Script para criar apresentação sobre Diagrama de Sequência UML
Curso: Modelagem de Sistemas - SENAI
Autor: Gerado com auxílio de IA
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.util import Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap

# Cores SENAI (em RGB hex)
AZUL_ESCURO = (0x3D, 0x4F, 0x5F)
AMARELO = (0xE8, 0xD4, 0x4D)
AZUL_CLARO = (0x8B, 0xB8, 0xD0)
VERDE = (0x4A, 0x7C, 0x6F)
BRANCO = (0xFF, 0xFF, 0xFF)
PRETO = (0x00, 0x00, 0x00)
AZUL_CARD = (0x2A, 0x3A, 0x4A)

def set_fill_color(shape, rgb_tuple):
    """Define cor de preenchimento"""
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(rgb_tuple)

def set_font_color(run_or_para, rgb_tuple):
    """Define cor da fonte"""
    run_or_para.font.color.rgb = _rgb(rgb_tuple)

def set_line_color(shape, rgb_tuple):
    """Define cor da linha"""
    shape.line.color.rgb = _rgb(rgb_tuple)

def _rgb(rgb_tuple):
    """Converte tupla RGB para objeto RGBColor"""
    from pptx.dml.color import RGBColor
    return RGBColor(rgb_tuple[0], rgb_tuple[1], rgb_tuple[2])

def criar_slide_titulo(prs, titulo, subtitulo=""):
    """Cria slide de título"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Layout em branco
    
    # Fundo
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = _rgb(AZUL_ESCURO)
    background.line.fill.background()
    
    # Título
    titulo_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = titulo_box.text_frame
    p = tf.paragraphs[0]
    p.text = titulo
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = _rgb(BRANCO)
    p.alignment = PP_ALIGN.CENTER
    
    if subtitulo:
        # Subtítulo em caixa amarela
        sub_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2), Inches(4), Inches(6), Inches(0.8))
        sub_box.fill.solid()
        sub_box.fill.fore_color.rgb = _rgb(AMARELO)
        sub_box.line.fill.background()
        
        sub_text = slide.shapes.add_textbox(Inches(2), Inches(4.1), Inches(6), Inches(0.6))
        tf = sub_text.text_frame
        p = tf.paragraphs[0]
        p.text = subtitulo
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = _rgb(PRETO)
        p.alignment = PP_ALIGN.CENTER
    
    return slide

def criar_slide_conteudo(prs, titulo, conteudo_lista, destaque=""):
    """Cria slide com título e lista de conteúdo"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Fundo
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = _rgb(AZUL_ESCURO)
    background.line.fill.background()
    
    # Barra de título
    barra = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(0.3), prs.slide_width, Inches(0.9))
    barra.fill.solid()
    barra.fill.fore_color.rgb = _rgb(VERDE)
    barra.line.fill.background()
    
    # Título
    titulo_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
    tf = titulo_box.text_frame
    p = tf.paragraphs[0]
    p.text = titulo
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = _rgb(BRANCO)
    
    # Conteúdo
    y_pos = 1.5
    for item in conteudo_lista:
        item_box = slide.shapes.add_textbox(Inches(0.7), Inches(y_pos), Inches(8.5), Inches(0.6))
        tf = item_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"• {item}"
        p.font.size = Pt(20)
        p.font.color.rgb = _rgb(BRANCO)
        y_pos += 0.55
    
    # Destaque (se houver)
    if destaque:
        dest_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(y_pos + 0.3), Inches(9), Inches(0.7))
        dest_box.fill.solid()
        dest_box.fill.fore_color.rgb = _rgb(AMARELO)
        dest_box.line.fill.background()
        
        dest_text = slide.shapes.add_textbox(Inches(0.6), Inches(y_pos + 0.4), Inches(8.8), Inches(0.5))
        tf = dest_text.text_frame
        p = tf.paragraphs[0]
        p.text = destaque
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = _rgb(PRETO)
        p.alignment = PP_ALIGN.CENTER
    
    return slide

def criar_slide_conceito(prs, titulo, definicao, exemplo_dia, exemplo_prog):
    """Cria slide de conceito com 2 exemplos"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Fundo
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = _rgb(AZUL_ESCURO)
    background.line.fill.background()
    
    # Barra de título
    barra = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(0.2), prs.slide_width, Inches(0.8))
    barra.fill.solid()
    barra.fill.fore_color.rgb = _rgb(VERDE)
    barra.line.fill.background()
    
    # Título
    titulo_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = titulo_box.text_frame
    p = tf.paragraphs[0]
    p.text = titulo
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = _rgb(BRANCO)
    
    # Definição em caixa amarela
    def_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.3), Inches(1.2), Inches(9.4), Inches(1))
    def_box.fill.solid()
    def_box.fill.fore_color.rgb = _rgb(AMARELO)
    def_box.line.fill.background()
    
    def_text = slide.shapes.add_textbox(Inches(0.4), Inches(1.3), Inches(9.2), Inches(0.8))
    tf = def_text.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = definicao
    p.font.size = Pt(18)
    p.font.color.rgb = _rgb(PRETO)
    
    # Exemplo 1 - Dia a dia
    ex1_titulo = slide.shapes.add_textbox(Inches(0.3), Inches(2.4), Inches(4.5), Inches(0.4))
    tf = ex1_titulo.text_frame
    p = tf.paragraphs[0]
    p.text = "🏠 Exemplo do Dia a Dia:"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = _rgb(AZUL_CLARO)
    
    ex1_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(2.8), Inches(4.5), Inches(1.5))
    ex1_box.fill.solid()
    ex1_box.fill.fore_color.rgb = _rgb(AZUL_CARD)
    ex1_box.line.color.rgb = _rgb(AZUL_CLARO)
    
    ex1_text = slide.shapes.add_textbox(Inches(0.4), Inches(2.9), Inches(4.3), Inches(1.3))
    tf = ex1_text.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = exemplo_dia
    p.font.size = Pt(14)
    p.font.color.rgb = _rgb(BRANCO)
    
    # Exemplo 2 - Programador
    ex2_titulo = slide.shapes.add_textbox(Inches(5.2), Inches(2.4), Inches(4.5), Inches(0.4))
    tf = ex2_titulo.text_frame
    p = tf.paragraphs[0]
    p.text = "💻 Exemplo do Programador:"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = _rgb(AMARELO)
    
    ex2_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.2), Inches(2.8), Inches(4.5), Inches(1.5))
    ex2_box.fill.solid()
    ex2_box.fill.fore_color.rgb = _rgb(AZUL_CARD)
    ex2_box.line.color.rgb = _rgb(AMARELO)
    
    ex2_text = slide.shapes.add_textbox(Inches(5.3), Inches(2.9), Inches(4.3), Inches(1.3))
    tf = ex2_text.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = exemplo_prog
    p.font.size = Pt(14)
    p.font.color.rgb = _rgb(BRANCO)
    
    return slide

def criar_slide_notacao(prs, titulo, simbolo_desc, como_desenhar):
    """Cria slide explicando notação visual"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Fundo
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = _rgb(AZUL_ESCURO)
    background.line.fill.background()
    
    # Barra de título
    barra = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(0.2), prs.slide_width, Inches(0.7))
    barra.fill.solid()
    barra.fill.fore_color.rgb = _rgb(VERDE)
    barra.line.fill.background()
    
    # Título
    titulo_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.6))
    tf = titulo_box.text_frame
    p = tf.paragraphs[0]
    p.text = titulo
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = _rgb(BRANCO)
    
    # Símbolo/Descrição
    simb_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.3), Inches(1.1), Inches(9.4), Inches(0.8))
    simb_box.fill.solid()
    simb_box.fill.fore_color.rgb = _rgb(AMARELO)
    simb_box.line.fill.background()
    
    simb_text = slide.shapes.add_textbox(Inches(0.4), Inches(1.2), Inches(9.2), Inches(0.6))
    tf = simb_text.text_frame
    p = tf.paragraphs[0]
    p.text = simbolo_desc
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = _rgb(PRETO)
    p.alignment = PP_ALIGN.CENTER
    
    # Como desenhar
    como_titulo = slide.shapes.add_textbox(Inches(0.3), Inches(2.1), Inches(9), Inches(0.4))
    tf = como_titulo.text_frame
    p = tf.paragraphs[0]
    p.text = "📐 Como Representar:"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = _rgb(AZUL_CLARO)
    
    como_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(2.5), Inches(9.4), Inches(2))
    como_box.fill.solid()
    como_box.fill.fore_color.rgb = _rgb(AZUL_CARD)
    como_box.line.color.rgb = _rgb(AZUL_CLARO)
    
    como_text = slide.shapes.add_textbox(Inches(0.5), Inches(2.6), Inches(9), Inches(1.8))
    tf = como_text.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = como_desenhar
    p.font.size = Pt(16)
    p.font.color.rgb = _rgb(BRANCO)
    
    return slide

# Criar apresentação
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)

# ========== SLIDE 1: CAPA ==========
criar_slide_titulo(prs, "Modelagem de Sistemas", "UML: Diagrama de Sequência")

# ========== SLIDE 2: O QUE É ==========
criar_slide_conteudo(prs, "O que é o Diagrama de Sequência?", [
    "É um diagrama COMPORTAMENTAL da UML",
    "Mostra a ORDEM (sequência) das interações entre objetos",
    "Representa eventos e mensagens ao longo do TEMPO",
    "Baseado no Diagrama de Casos de Uso e Classes",
    "Normalmente: 1 Diagrama de Sequência para CADA Caso de Uso"
], "📌 Foco: COMO os objetos conversam entre si, passo a passo")

# ========== SLIDE 3: ELEMENTOS BÁSICOS ==========
criar_slide_conteudo(prs, "Elementos Básicos do Diagrama", [
    "ATOR: Quem inicia a interação (pessoa, sistema externo)",
    "OBJETO: Instância de uma classe (ex: :Pedido, :Cliente)",
    "LINHA DE VIDA: Linha tracejada vertical (existência do objeto)",
    "MENSAGEM: Seta horizontal (comunicação entre objetos)",
    "ATIVAÇÃO: Retângulo fino (objeto processando)"
], "⏰ O tempo flui de CIMA para BAIXO no diagrama")

# ========== SLIDE 4: OBJETO - EXPLICAÇÃO ==========
criar_slide_conceito(prs, 
    "Elemento: OBJETO (Lifeline)",
    "O objeto é representado por um RETÂNGULO no topo do diagrama. Dentro dele: nome_da_instância : NomeDaClasse. Exemplo: pedido1 : Pedido ou apenas :Pedido (anônimo)",
    "Pense em uma PESSOA em uma reunião. O retângulo é o 'crachá' dela com o nome. A linha abaixo mostra que ela está presente na reunião.",
    "É como uma VARIÁVEL no código: Pedido pedido1 = new Pedido(); O retângulo representa essa variável 'viva' na memória."
)

# ========== SLIDE 5: COMO DESENHAR OBJETO ==========
criar_slide_notacao(prs,
    "Como Desenhar um OBJETO",
    "📦 Retângulo com texto sublinhado: nomeobjeto : Classe",
    """1. Desenhe um RETÂNGULO no topo do diagrama
2. Dentro, escreva: nome : Classe (ex: car1 : Carrinho)
3. O nome pode ser omitido: apenas :Classe
4. O texto deve ser SUBLINHADO (indica instância)
5. Abaixo, sai uma LINHA TRACEJADA vertical (linha de vida)

Exemplos válidos:
• pesfis1 : Pessoa_Fisica
• :Pedido (objeto anônimo)
• comum1 : Conta_Comum"""
)

# ========== SLIDE 6: ATOR ==========
criar_slide_conceito(prs,
    "Elemento: ATOR",
    "O ator é representado pelo 'boneco palito' (stickman). É quem INICIA a interação com o sistema. Pode ser uma pessoa, outro sistema, ou dispositivo externo.",
    "É como o CLIENTE que chega no restaurante e faz o pedido. Ele dispara a ação, mas a cozinha (sistema) faz o trabalho.",
    "É quem clica no botão, envia requisição HTTP, ou dispara um evento. O ator NÃO faz parte do código, mas INTERAGE com ele."
)

# ========== SLIDE 7: LINHA DE VIDA ==========
criar_slide_conceito(prs,
    "Elemento: LINHA DE VIDA (Lifeline)",
    "Linha vertical TRACEJADA que desce do objeto. Representa o TEMPO de existência do objeto durante a interação. Pode ser interrompida por um X quando o objeto é destruído.",
    "Como a 'duração' de uma ligação telefônica. Enquanto a linha existe, a pessoa está na chamada. Quando a linha termina, a chamada acabou.",
    "Enquanto a linha existe, o objeto está na MEMÓRIA. É como o escopo de uma variável - existe enquanto o método/bloco está executando."
)

# ========== SLIDE 8: MENSAGEM SIMPLES ==========
criar_slide_conceito(prs,
    "Elemento: MENSAGEM (Chamada de Método)",
    "Seta HORIZONTAL que vai de uma linha de vida para outra. Representa uma COMUNICAÇÃO: chamada de método, envio de dados, ou ação. O nome da mensagem fica sobre a seta.",
    "Cliente → Garçom: 'Por favor, um café'. Essa frase é a MENSAGEM. O garçom recebe e processa.",
    "É como chamar uma FUNÇÃO: objeto.metodo(). A seta mostra QUEM chama e QUEM recebe. Ex: sistema.validar(cpf)"
)

print("Slides 1-8 criados com sucesso!")

# ========== SLIDE 9: MENSAGEM ENTRE ATORES ==========
criar_slide_conceito(prs,
    "Mensagem entre ATORES",
    "Quando dois ATORES se comunicam, a seta é SÓLIDA conectando suas linhas de vida. NÃO é chamada de método - é comunicação de alto nível (fala, clique, gesto).",
    "Exemplo 1: Cliente → Funcionário: 'Quero abrir uma conta'. Exemplo 2: Paciente → Recepcionista: 'Preciso marcar consulta'.",
    "Exemplo 1: Usuário → Administrador: 'Reset minha senha'. Exemplo 2: Cliente → Suporte: 'Sistema travou'. É comunicação FORA do código."
)

# ========== SLIDE 10: MENSAGEM ENTRE OBJETOS ==========
criar_slide_conceito(prs,
    "Mensagem entre OBJETOS",
    "Seta SÓLIDA (cheia) de uma linha de vida para outra. Representa chamada de MÉTODO. O texto é o nome do método com parâmetros: metodo(param).",
    "Exemplo 1: Caixa → Calculadora: somar(10, 5). Exemplo 2: Atendente → Estoque: verificar('Pastel').",
    "Exemplo 1: controlador.validarCPF('123'). Exemplo 2: bd.consultar(SELECT * FROM users). É código chamando código!"
)

# ========== SLIDE 11: CRIAÇÃO DE OBJETO ==========
criar_slide_conceito(prs,
    "CRIAÇÃO de Objeto",
    "Quando um objeto cria OUTRO objeto durante a execução. A seta aponta para o RETÂNGULO do novo objeto (não para a linha). O novo objeto aparece mais abaixo no diagrama.",
    "Exemplo 1: Garçom cria uma COMANDA quando cliente faz pedido. Exemplo 2: Caixa cria NOTA FISCAL quando fecha a venda.",
    "Exemplo 1: Pedido p = new Pedido(). Exemplo 2: Item item = new Item_Pedido(). É o 'new' criando instância!"
)

# ========== SLIDE 12: COMO DESENHAR CRIAÇÃO ==========
criar_slide_notacao(prs,
    "Como Desenhar CRIAÇÃO de Objeto",
    "📦 Seta aponta para o RETÂNGULO do novo objeto (não para linha)",
    """1. O novo objeto NÃO aparece no topo (não existe ainda)
2. Desenhe a seta com <<create>> ou nome do construtor
3. A seta aponta DIRETAMENTE para o retângulo do objeto
4. O objeto é desenhado mais ABAIXO, onde é criado
5. A partir daí, sua linha de vida começa

Notação da seta: ─────────────> □ :NovoObjeto
                  new() ou <<create>>"""
)

# ========== SLIDE 13: DESTRUIÇÃO DE OBJETO ==========
criar_slide_conceito(prs,
    "DESTRUIÇÃO de Objeto",
    "Quando um objeto deixa de existir. A linha de vida termina com um X grande. Representa remoção da memória ou fim do escopo.",
    "Exemplo 1: Você TERMINA uma ligação telefônica (☎️→ X). Exemplo 2: Restaurante ENCERRA a mesa após pagamento.",
    "Exemplo 1: objeto = null (garbage collection). Exemplo 2: conexao.close(). Exemplo: session.invalidate() no logout."
)

# ========== SLIDE 14: COMO DESENHAR DESTRUIÇÃO ==========
criar_slide_notacao(prs,
    "Como Desenhar DESTRUIÇÃO de Objeto",
    "❌ Um X grande no final da linha de vida",
    """1. Desenhe uma seta com <<destroy>> ou método de exclusão
2. No final da linha de vida do objeto destruído: X
3. A linha de vida NÃO continua após o X
4. Significa que o objeto saiu da memória

Representação:
    :Objeto
       │
       │  ←── excluir()
       X (fim da existência)"""
)

# ========== SLIDE 15: MENSAGEM DE RETORNO ==========
criar_slide_conceito(prs,
    "MENSAGEM de Retorno",
    "Seta TRACEJADA (pontilhada) que volta para quem chamou. Representa a RESPOSTA após processar um método. Pode incluir o valor retornado.",
    "Exemplo 1: Garçom → Cliente: 'Seu pedido está pronto' (resposta). Exemplo 2: Porteiro → Morador: 'Sua encomenda chegou'.",
    "Exemplo 1: return true; Exemplo 2: return cliente.getNome(); A seta mostra O QUE volta: dados, confirmação, erro."
)

# ========== SLIDE 16: COMO DESENHAR RETORNO ==========
criar_slide_notacao(prs,
    "Como Desenhar MENSAGEM de Retorno",
    "◀- - - - - Seta TRACEJADA apontando de volta",
    """1. Use linha TRACEJADA (- - - - -)
2. Seta aponta para QUEM fez a chamada original
3. Texto opcional: nome_retorno : tipo
4. Pode ser IMPLÍCITO (não obrigatório desenhar)

Exemplo completo:
    :Sistema ──────────────> :BD
              consultar()
    :Sistema <─ ─ ─ ─ ─ ─ ─  :BD
              dados : String"""
)

print("Slides 9-16 criados com sucesso!")

# ========== SLIDE 17: AUTOCHAMADA ==========
criar_slide_conceito(prs,
    "AUTOCHAMADA (Self-Call)",
    "Quando um objeto chama um método DELE MESMO. A seta SAI da linha de vida e VOLTA para a mesma linha, fazendo uma curva.",
    "Exemplo 1: Você faz um 'checklist mental' antes de sair de casa. Exemplo 2: Calculadora 'confere' o resultado antes de mostrar.",
    "Exemplo 1: this.validarDados(). Exemplo 2: self.calcularTotal(). É quando o objeto processa algo internamente."
)

# ========== SLIDE 18: COMO DESENHAR AUTOCHAMADA ==========
criar_slide_notacao(prs,
    "Como Desenhar AUTOCHAMADA",
    "↩️ Seta que sai e volta para a MESMA linha de vida",
    """1. A seta SAI do retângulo de ativação
2. Faz uma CURVA para a direita
3. VOLTA para o mesmo objeto
4. Cria um retângulo de ativação SOBREPOSTO

    :Objeto
       ┌─┐
       │ │──┐ metodoInterno()
       │ │  │
       │ │──┘
       │ │
       └─┘"""
)

# ========== SLIDE 19: RESTRIÇÕES DE TEMPO ==========
criar_slide_conceito(prs,
    "RESTRIÇÕES de Tempo",
    "Quando uma ação DEMORA ou tem PRAZO definido. A seta é desenhada na DIAGONAL (inclinada), indicando passagem de tempo. Pode ter nota com o tempo.",
    "Exemplo 1: 'Aguarde 30 minutos para a pizza ficar pronta'. Exemplo 2: 'Seu pedido será cancelado em 15 min se não confirmar'.",
    "Exemplo 1: Thread.sleep(5000) - espera 5 segundos. Exemplo 2: setTimeout() em JavaScript. Exemplo: API com timeout de 30s."
)

# ========== SLIDE 20: COMO DESENHAR RESTRIÇÃO TEMPO ==========
criar_slide_notacao(prs,
    "Como Desenhar RESTRIÇÕES de Tempo",
    "⏱️ Seta DIAGONAL + nota com tempo entre chaves { }",
    """1. Desenhe a seta em DIAGONAL (inclinada para baixo)
2. Adicione nota: {tempo de espera}
3. A inclinação mostra que TEMPO passou

    :Sistema                :Pedido
       │                       │
       │\\                      │
       │ \\ cancelar()         │
       │  \\ {30 min de espera}│
       │   \\________________> X

O pedido só é cancelado APÓS 30 minutos."""
)

# ========== SLIDE 21: FRAGMENTOS - INTRODUÇÃO ==========
criar_slide_conteudo(prs, "FRAGMENTOS de Interação", [
    "São 'caixas' que agrupam partes do diagrama",
    "Definem comportamentos especiais: condição, repetição, etc.",
    "Etiqueta no canto superior esquerdo indica o tipo",
    "Principais tipos: ref, opt, alt, loop, par",
    "Funcionam como estruturas de controle do código"
], "📦 Pense como: IF, ELSE, FOR, WHILE... mas visual!")

# ========== SLIDE 22: FRAGMENTO REF ==========
criar_slide_conceito(prs,
    "Fragmento: REF (Referência)",
    "Um 'LINK' para outro diagrama. Evita repetir desenhos. O retângulo atravessa as linhas de vida envolvidas. Texto central = nome do diagrama referenciado.",
    "Exemplo 1: Manual de instruções que diz 'Veja página 45'. Exemplo 2: Receita que diz 'Para o molho, veja receita X'.",
    "Exemplo 1: Chamar função de outro arquivo: import login; login.autenticar(). Exemplo 2: Módulo separado que é reutilizado."
)

# ========== SLIDE 23: FRAGMENTO OPT ==========
criar_slide_conceito(prs,
    "Fragmento: OPT (Opcional)",
    "É o 'IF sozinho' - sem ELSE. A caixa só executa SE a condição for verdadeira. Se falsa, pula e continua. Condição entre colchetes [condição].",
    "Exemplo 1: 'SE estiver chovendo, leve guarda-chuva'. Exemplo 2: 'SE cliente for VIP, dar desconto'.",
    "Exemplo 1: if (usuario.isPremium()) { darBrinde(); }. Exemplo 2: if (carrinho.valor > 100) { freteGratis(); }"
)

# ========== SLIDE 24: FRAGMENTO ALT ==========
criar_slide_conceito(prs,
    "Fragmento: ALT (Alternativo)",
    "É o 'IF / ELSE'. Caixa dividida ao meio por linha tracejada. Parte de cima: condição verdadeira. Parte de baixo: condição falsa (else).",
    "Exemplo 1: 'SE aprovado, comemorar. SENÃO, estudar mais'. Exemplo 2: 'SE tem saldo, comprar. SENÃO, pedir empréstimo'.",
    "if (senha.valida()) { login(); } else { erroSenha(); }. Só UM caminho executa!"
)

print("Slides 17-24 criados com sucesso!")

# ========== SLIDE 25: FRAGMENTO LOOP ==========
criar_slide_conceito(prs,
    "Fragmento: LOOP (Repetição)",
    "É o 'FOR / WHILE'. As ações dentro da caixa se REPETEM enquanto condição válida. Condição entre colchetes. Etiqueta 'loop' no canto.",
    "Exemplo 1: 'PARA CADA item do carrinho, calcular preço'. Exemplo 2: 'ENQUANTO tiver fome, comer mais um pastel'.",
    "for (Item i : carrinho) { total += i.preco; }. Ou: while (fila.hasNext()) { processar(); }"
)

# ========== SLIDE 26: COMO DESENHAR FRAGMENTOS ==========
criar_slide_notacao(prs,
    "Como Desenhar FRAGMENTOS",
    "📦 Retângulo com etiqueta (aba) no canto superior esquerdo",
    """1. Desenhe um RETÂNGULO envolvendo as mensagens
2. No canto superior esquerdo: 'aba' com tipo (opt, alt, loop)
3. Abaixo da aba: [condição] entre colchetes
4. Para ALT: linha tracejada horizontal divide as opções

    ┌─ opt ────────────────────────┐
    │ [se cliente VIP]             │
    │   :Sistema ──> :Desconto     │
    │              aplicar()       │
    └──────────────────────────────┘"""
)

# ========== SLIDE 27: MENSAGEM PERDIDA/ENCONTRADA ==========
criar_slide_conceito(prs,
    "Mensagem PERDIDA e ENCONTRADA",
    "PERDIDA: não chegou ao destino (destino fora do diagrama). ENCONTRADA: veio de origem desconhecida. Representadas com círculo preto no início/fim.",
    "Exemplo PERDIDA: Você envia email mas não sabe se chegou. Exemplo ENCONTRADA: Recebe notificação sem saber quem enviou.",
    "PERDIDA: API externa que não responde. ENCONTRADA: Webhook recebido de serviço externo não modelado."
)

# ========== SLIDE 28: RESUMO VISUAL ==========
criar_slide_conteudo(prs, "📋 Resumo: Elementos Principais", [
    "RETÂNGULO sublinhado = Objeto (:Classe)",
    "BONECO PALITO = Ator (usuário/sistema externo)",
    "LINHA TRACEJADA vertical = Linha de Vida",
    "SETA SÓLIDA horizontal = Mensagem/Chamada",
    "SETA TRACEJADA = Retorno",
    "RETÂNGULO fino = Ativação (processando)",
    "X = Destruição do objeto",
    "CAIXA com aba = Fragmento (opt, alt, loop, ref)"
], "")

# ========== SLIDE 29: DICAS PRÁTICAS ==========
criar_slide_conteudo(prs, "💡 Dicas para Fazer Bons Diagramas", [
    "Comece pelo ATOR que inicia a ação",
    "Tempo flui de CIMA para BAIXO",
    "Use nomes de MÉTODOS reais (verbo + substantivo)",
    "Não esqueça dos RETORNOS importantes",
    "Use FRAGMENTOS para condições e loops",
    "Mantenha o diagrama SIMPLES e legível",
    "Um diagrama por CASO DE USO"
], "🎯 Objetivo: Qualquer desenvolvedor deve entender o fluxo!")

# ========== SLIDE 30: ENCERRAMENTO ==========
criar_slide_titulo(prs, "Pratique!", "Agora é sua vez de criar diagramas")

print("Slides 25-30 criados com sucesso!")
print("=" * 50)
print("APRESENTAÇÃO COMPLETA!")

# Salvar apresentação
output_path = r"c:\Users\nadson\Desktop\modelagem-de-sistema\Planos de Aula\Encontro 4 - 05DEZ - Sequência e Atividades\Diagrama_Sequencia_NOVO.pptx"
prs.save(output_path)
print(f"Apresentação salva em: {output_path}")
